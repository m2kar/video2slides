# Copyright (c) 2024 m2kar
# Created: 2024/12/01
# Author: m2kar <m2kar.cn#gmail.com>
# License: MIT
# File: main.py
# Description: main.py for video2slides, convert video to pptx slides

# video2slides/main.py
import os
import sys
import shutil
import typer
from typing_extensions import Annotated
from typing import Optional

import datetime
import re
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tqdm
import signal
import atexit

DEBUG=bool(os.getenv("DEBUG",False))

class Config:
    skip_interval: int = 0.15 # 跳过的时间间隔 (秒)
    skip_frames: int = 5 # 默认跳过的帧数，在无法获取视频帧率的情况下使用
    
    threshold: float = 0.1  # 提高阈值以减少误切
    min_scene_len: float = 1.0  # 最小场景时间长度(s)
    
    screenshot_position: float = 0.95 # 截图在场景中的位置 (默认的时间位置比例)
    min_screenshot_margin: float = 0.05 # 截图与场景结束的最小间隔 (秒)
    max_screenshot_margin: float = 0.5 # 截图与场景结束的最大间隔 (秒)

    
def make_pptx(output:str,img_list:list,note_list:list=None):
    prs=Presentation()
    
    with open(img_list[0], 'rb') as f:
        im = Image.open(f)
    dpi=im.info.get('dpi',(96,96))
    width_px, height_px = im.size
    width_in = width_px / dpi[0]
    height_in = height_px / dpi[1]
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    
    blank_slide_layout = prs.slide_layouts[6]
    for i,img in enumerate(img_list):
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = 0
        pic = slide.shapes.add_picture(img, left, top, width=prs.slide_width, height=prs.slide_height)
        
        # add note
        if note_list is not None:
            note = note_list[i]
            slide.notes_slide.notes_text_frame.text = note
        
    prs.save(output)
    
# Scene Number,Start Frame,Start Timecode,Start Time (seconds),End Frame,End Timecode,End Time (seconds),Length (frames),Length (timecode),Length (seconds)
# 1,1,00:00:00.000,0.000,174,00:00:06.960,6.960,174,00:00:06.960,6.960
# 7,16405,00:10:56.160,656.160,18024,00:12:00.960,720.960,1620,00:01:04.800,64.800
def time_info(scene_info_csv:str):
    scene_infos=load_scene_information(scene_info_csv)
    time_notes=[ f"Time Info: scene {scene_details[0]}, {scene_details[2]}-{scene_details[5]} ({scene_details[9]}s)"
                for scene_details in scene_infos]
    return time_notes

def load_scene_information(scene_info_csv:str):
    with open(scene_info_csv) as f:
        f.readline()
        f.readline()
        lines=f.readlines()
    scene_infos=[l.strip().split(",") for l in lines]
    return scene_infos
 
def seconds_to_timecode(seconds):
    return f"{int(seconds//3600):02d}:{int(seconds%3600//60):02d}:{seconds%60:06.3f}"

def srt_info(scene_info_csv:str,video_path:str):
    import whisper

    # 加载模型
    model = whisper.load_model("turbo")
    # 转录音轨
    print("[*] Converting video to text...")
    result = model.transcribe(video_path,fp16=False,verbose=False)
    
    segments=result["segments"]
    print("[*] Convert video to srt Done.")
        
    scene_infos=load_scene_information(scene_info_csv)
    time_notes=[ f"Time Info: scene {scene_details[0]}, {scene_details[2]}-{scene_details[5]} ({scene_details[9]}s)"
                for scene_details in scene_infos]
    
    srt_notes = [""] * len(scene_infos)
    for seg in segments:
        seg_start = seg["start"]
        seg_end = seg["end"]
        seg_text = seg["text"]
        
        max_intersection = 0
        assigned_scene_index = -1
        
        for idx, scene in enumerate(scene_infos):
            scene_start = float(scene[3])  # 假设 start_seconds 在索引3
            scene_end = float(scene[6])    # 假设 end_seconds 在索引6
            
            # 计算交集时间
            overlap_start = max(seg_start, scene_start)
            overlap_end = min(seg_end, scene_end)
            overlap = max(0, overlap_end - overlap_start)
            
            if overlap > max_intersection:
                max_intersection = overlap
                assigned_scene_index = idx
        
        if assigned_scene_index != -1:
            # 添加字幕到对应的场景备注
            if srt_notes[assigned_scene_index]:
                srt_notes[assigned_scene_index] += "\n" + seg_text
            else:
                srt_notes[assigned_scene_index] = seg_text
    
    for i,scene in enumerate(scene_infos):
        srt_notes[i]+="\n\n" + time_notes[i]
    return srt_notes

def screenshots(scene_info_csv:str,video_path:str,tmp_dir:str):
    """screenshots from video, per scene per screenshot"""
    scene_infos=load_scene_information(scene_info_csv)
    imgs_list=[]
    img_basename_template="scene{scene_number:0"+str(len(str(len(scene_infos))))+"d}.jpg"
    for scene in tqdm.tqdm(scene_infos,desc="screenshotting"):
        scene_number=int(scene[0])
        start_seconds=float(scene[3])
        end_seconds=float(scene[6])
        
        default_screenshot_seconds=start_seconds+Config.screenshot_position*(end_seconds-start_seconds)
        screenshot_margin=max(Config.min_screenshot_margin,min(Config.max_screenshot_margin,end_seconds-default_screenshot_seconds))
        screenshot_seconds=end_seconds-screenshot_margin
        screenshot_seconds=round(screenshot_seconds,3)
        
        dst_path=os.path.join(tmp_dir,img_basename_template.format(scene_number=scene_number))
        screenshot_cmd=f"ffmpeg -ss {screenshot_seconds} -i '{video_path}' -frames:v 1 -q:v 2 -y '{dst_path}' "
        if not DEBUG: screenshot_cmd+=" > /dev/null 2>&1"
        if DEBUG: print(screenshot_cmd)
        os.system(screenshot_cmd)
        imgs_list.append(dst_path)
    return imgs_list

def get_video_fps(video_path:str):
    # 使用 ffprobe 获取视频信息

    cmd = f"ffprobe -v error -select_streams v:0 -show_entries stream=avg_frame_rate -of default=noprint_wrappers=1:nokey=1 '{video_path}'"
    result = os.popen(cmd).read().strip()
    
    # 解析帧率信息
    match = re.match(r'(\d+)/(\d+)', result)
    if match:
        num, denom = map(int, match.groups())
        fps = num / denom
        return fps
    else:
        raise ValueError("无法解析视频的帧率信息")

def get_skip_frames(video_path:str):
    try:
        # 获取视频的帧率
        fps = get_video_fps(video_path)-1
        fps = max(0, fps)
        
        # 计算跳过的帧数
        skip_frames = int(fps * Config.skip_interval)
        if DEBUG: print(f"[*] 视频帧率: {fps}，跳过帧数: {skip_frames}")
        return skip_frames
    except Exception as e:
        print(f"[!] 获取视频帧率失败，返回默认跳过帧数 skip_frame={Config.skip_frames}: {e}")
        return Config.skip_frames

def detect_scenes(video_path:str,tmp_dir:str):
    Config.skip_frames=get_skip_frames(video_path)
    """detect scenes from video"""
    scenedetect_cmd = (
        f" scenedetect -i '{video_path}' -o '{tmp_dir}' "
        f" -fs {Config.skip_frames} "
        f" detect-hash --size 16  --threshold {Config.threshold}  --lowpass 2 "
        f" --min-scene-len {Config.min_scene_len} "
        f" list-scenes --filename scenes_info.csv "
    )

    if not DEBUG: scenedetect_cmd+=" 2>&1"
    if DEBUG: print("[*]"+scenedetect_cmd)
    os.system(scenedetect_cmd)
    return os.path.join(tmp_dir,"scenes_info.csv")

def register_signal_handlers(local_vars):
    
    def cleanup_temp_dir():
        tmp_dir=local_vars.get("tmp_dir")
        if DEBUG:
            print(f"Not clean tmp_dir for debug:  {tmp_dir}")
        else:
            if tmp_dir and os.path.exists(tmp_dir):
                print(f"[*] Cleaning up temporary directory: {tmp_dir}")
                shutil.rmtree(tmp_dir)
    
    def signal_handler(sig, frame):
        print(f"[*] Signal {signum} received, cleaning up...")
        cleanup_temp_dir()
        raise SystemExit(1)
        
    atexit.register(cleanup_temp_dir)
    signal.signal(signal.SIGINT, signal_handler)
    signal.signal(signal.SIGTERM, signal_handler)

app = typer.Typer()

@app.command()
def main(
    video: Annotated[str, typer.Argument(..., help="输入视频文件路径")],
    output: Annotated[Optional[str], typer.Argument(...,help="输出pptx文件路径")] =None,
    srt: bool = typer.Option(False, help="是否生成字幕"),
    
    # srt: Annotated[Optional[bool], typer.Option(False, help="是否生成字幕")]=False
    ):
    """Convert video to pptx slides"""
    
    # 捕获 Ctrl+C 信号
    
    timestamp=datetime.datetime.now().strftime('%Y-%m-%d-%H-%M-%S')
    fname=os.path.basename(video)
    if output is None:
        output = os.path.splitext(fname)[0]+".pptx"
        if os.path.exists(output):
            output = os.path.splitext(fname)[0]+f"-{timestamp}.pptx"
    tmp_dir=f"./tmp-v2s-{timestamp}-{fname}"
    os.makedirs(tmp_dir,exist_ok=True)

    register_signal_handlers(locals())
    
    scene_info_csv=detect_scenes(video,tmp_dir)

    # get the note list
    if srt:
        notes_list=srt_info(scene_info_csv,video)
    else:
        time_notes=time_info(scene_info_csv)
        notes_list=time_notes
    
    # screenshot from scenes
    imgs_list = screenshots(scene_info_csv,video_path=video,tmp_dir=tmp_dir)

    # convert to pptx
    make_pptx(output,imgs_list,notes_list)
    
    print(f"已生成 {output}")
    
if __name__ == "__main__":
    app()